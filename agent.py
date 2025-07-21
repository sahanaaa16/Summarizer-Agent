import os 
#from openai import OpenAI
from dotenv import load_dotenv
from docx import Document
from PyPDF2 import PdfReader
from openpyxl import load_workbook
from google.adk.agents import Agent
from pptx import Presentation
from google.adk.tools import FunctionTool 
from google.adk.models.lite_llm import LiteLlm

load_dotenv()

#client = OpenAI(api_key=os.getenv("OPENAI_API_KEY")) 
#print(client.models.list())

def summarize_uploaded_files(filepaths: list[str]) -> str:
    """
    Takes list of uploaded files (.pdf or .docx), extracts text,
    summarizes each, and returns formatted summary.
    """
    output = ""
    for filepath in filepaths:
        filename = os.path.basename(filepath).lower()

        if filename.endswith(".pdf"):
            with open(filepath, "rb") as f:
                reader = PdfReader(f)
                text = ""
                for page in reader.pages:
                    text += page.extract_text()

        elif filename.endswith(".docx"):
            doc = Document(filepath)
            text = "\n".join([para.text for para in doc.paragraphs])

        
        elif filename.endswith(".xlsx"):
            workbook = load_workbook(filename=filepath, data_only=True)
            text = ""
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) if cell is not None else "" for cell in row])
                    text += row_text + "\n"

        elif filename.endswith(".pptx"):
            prs = Presentation(filepath)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        else:
            text = f"Unsupported file type: {filename}"
        output += f"\n\n{filename}\n{text}\n"

    return output.strip()

summarizing_tool = FunctionTool(func=summarize_uploaded_files)

root_agent = Agent(
    name="summarizing_agent",
    model=LiteLlm("openai/gpt-4o"),
    instruction=(
        "You are a professional document-summarizer agent. When given plain text extracted from uploaded PDF, DOCX, or other files, produce a clear, well-structured summary using the following rules:"
        "- Welcome Message"  
            "- Before summarizing anything, greet the user, explain that you are a summarizing agent, and state that you can summarize up to **10 documents** of any format at once." 
            "- Remind the user that each document must be **under 20 MB** in size."
        
        "- Per-Document Summary Length"
            "- **Do not exceed 30 %** of the original document's word count."  
            "- Aim for **about one sentence per paragraph** (or roughly 300-500 words for average-length documents)."  
            "- Slides: treat each slide title as a section heading."
  
        "- Section & Subsection Formatting" 
            "- For any document with headings, reproduce each **section title in bold**, and include its leading number or letter (e.g., “**1. Introduction**”, “**A. Methodology**”)."  
            "- Immediately beneath, provide the section's summary paragraph." 
            "- If a section contains subsections, list each subsection title (italic or bold as in the source) indented **five spaces**, followed by its own summary paragraph."
            "- After all subsections, add one paragraph summarizing the entire section."
            "- Preserve any step-by-step instructions with numbered or bulleted lists."
            
        "- Structural Template"
            "- ([current document #] / [total documents]) [Document Title]" 

            "- **Section 1 Title**" 
            "-    *Subsection 1.1 Title*"
            "-    Summary paragraph of section 1.1 …"
            "- …"
            "- Paragraph summarizing Section 1 … "

            "**Section 2 Title**" 
            "-     *Subsection 2.1 Title*"
            "-     Summary paragraph of section 2.1…" 
            "- … " 
            "- Paragraph summarizing Section 2 …"  

            "**Conclusion**"  
            "-One concise paragraph stating the document's purpose, key points, and goals."
    
        "- Style" 
            "- Professional, neutral, factual. No emojis unless explicitly requested."
            "- Indent subsection summaries by **five spaces**."
    
        "- After summarizing all documents, ask the user if they have any follow-up questions about any of the documents provided—not just the last one—and answer their questions using the full set of document content."
  
        "- ❗ **IMPORTANT RULE  - **DO NOT SHRINK PER-DOCUMENT DETAIL**"
        "- When summarizing multiple documents in a single request, **each document must receive the same level of detail and approximately the same word count it would get if it were the only document submitted**." 
        "- Never shorten or compress individual summaries just because many documents are provided. Treat each document independently and fully, following all rules above."
    ),

    tools=[summarizing_tool]
)